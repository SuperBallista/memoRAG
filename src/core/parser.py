"""문서 파서 모듈 - PDF, DOCX, HWPX, TXT, MD 지원"""
from pathlib import Path
from typing import List, Dict, Optional
import logging
from dataclasses import dataclass

logger = logging.getLogger(__name__)


@dataclass
class DocumentChunk:
    """문서 청크 데이터 클래스"""
    text: str
    metadata: Dict[str, any]
    page: Optional[int] = None
    section: Optional[str] = None


class DocumentParser:
    """다양한 문서 포맷을 파싱하는 클래스"""
    
    SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".hwpx", ".txt", ".md", ".pptx", ".xlsx"}
    
    def __init__(self, chunk_size: int = 512, chunk_overlap: int = 50):
        """
        Args:
            chunk_size: 청크 크기 (문자 단위)
            chunk_overlap: 청크 간 중복 (문자 단위)
        """
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
    
    def parse(self, file_path: Path) -> List[DocumentChunk]:
        """
        파일을 파싱하여 청크 리스트로 반환
        
        Args:
            file_path: 파싱할 파일 경로
            
        Returns:
            DocumentChunk 리스트
            
        Raises:
            ValueError: 지원하지 않는 파일 형식
            FileNotFoundError: 파일이 존재하지 않음
        """
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        extension = file_path.suffix.lower()
        if extension not in self.SUPPORTED_EXTENSIONS:
            raise ValueError(f"Unsupported file type: {extension}")
        
        logger.info(f"Parsing file: {file_path}")
        
        # 파일 타입별 파싱
        if extension == ".pdf":
            return self._parse_pdf(file_path)
        elif extension == ".docx":
            return self._parse_docx(file_path)
        elif extension == ".hwpx":
            return self._parse_hwpx(file_path)
        elif extension == ".pptx":
            return self._parse_pptx(file_path)
        elif extension == ".xlsx":
            return self._parse_xlsx(file_path)
        elif extension in {".txt", ".md"}:
            return self._parse_text(file_path)
        
        return []
    
    def _parse_pdf(self, file_path: Path) -> List[DocumentChunk]:
        """PDF 파일 파싱 (pypdf 사용)"""
        try:
            import pypdf
        except ImportError:
            raise ImportError("pypdf is required for PDF parsing. Install: pip install pypdf")
        
        chunks = []
        with open(file_path, "rb") as f:
            reader = pypdf.PdfReader(f)
            
            for page_num, page in enumerate(reader.pages, start=1):
                text = page.extract_text()
                if text.strip():
                    # 페이지 텍스트를 청크로 분할
                    page_chunks = self._split_text(text)
                    for chunk_text in page_chunks:
                        chunks.append(DocumentChunk(
                            text=chunk_text,
                            page=page_num,
                            metadata={
                                "file_path": str(file_path),
                                "file_name": file_path.name,
                                "file_type": "pdf",
                                "page": page_num,
                                "total_pages": len(reader.pages)
                            }
                        ))
        
        logger.info(f"Parsed {len(chunks)} chunks from PDF")
        return chunks
    
    def _parse_docx(self, file_path: Path) -> List[DocumentChunk]:
        """DOCX 파일 파싱 (python-docx 사용)"""
        try:
            from docx import Document
        except ImportError:
            raise ImportError("python-docx is required. Install: pip install python-docx")
        
        chunks = []
        doc = Document(file_path)
        
        # 모든 단락 텍스트 추출
        full_text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
        
        # 청크로 분할
        chunk_texts = self._split_text(full_text)
        for idx, chunk_text in enumerate(chunk_texts):
            chunks.append(DocumentChunk(
                text=chunk_text,
                metadata={
                    "file_path": str(file_path),
                    "file_name": file_path.name,
                    "file_type": "docx",
                    "chunk_index": idx
                }
            ))
        
        logger.info(f"Parsed {len(chunks)} chunks from DOCX")
        return chunks
    
    def _parse_hwpx(self, file_path: Path) -> List[DocumentChunk]:
        """HWPX 파일 파싱 (zip + XML 구조)"""
        try:
            from lxml import etree
            import zipfile
        except ImportError:
            raise ImportError("lxml is required. Install: pip install lxml")
        
        chunks = []
        
        try:
            with zipfile.ZipFile(file_path, 'r') as z:
                # HWPX는 ZIP 구조이며, Contents/section*.xml에 텍스트가 있음
                section_files = [name for name in z.namelist() if name.startswith('Contents/section')]
                
                full_text = ""
                for section_file in sorted(section_files):
                    with z.open(section_file) as f:
                        tree = etree.parse(f)
                        # 텍스트 노드 추출 (간단 버전)
                        text_nodes = tree.xpath("//text()")
                        section_text = " ".join([t.strip() for t in text_nodes if t.strip()])
                        full_text += section_text + "\n"
                
                # 청크로 분할
                chunk_texts = self._split_text(full_text)
                for idx, chunk_text in enumerate(chunk_texts):
                    chunks.append(DocumentChunk(
                        text=chunk_text,
                        metadata={
                            "file_path": str(file_path),
                            "file_name": file_path.name,
                            "file_type": "hwpx",
                            "chunk_index": idx
                        }
                    ))
        except Exception as e:
            logger.error(f"Error parsing HWPX: {e}")
            raise
        
        logger.info(f"Parsed {len(chunks)} chunks from HWPX")
        return chunks
    
    def _parse_pptx(self, file_path: Path) -> List[DocumentChunk]:
        """PPTX 파일 파싱 (python-pptx 사용)"""
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("python-pptx is required. Install: pip install python-pptx")
        
        chunks = []
        prs = Presentation(file_path)
        
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text = []
            
            # 제목
            if slide.shapes.title:
                slide_text.append(slide.shapes.title.text)
            
            # 본문 텍스트
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            # 슬라이드 텍스트 결합
            full_text = "\n".join(slide_text)
            
            if full_text.strip():
                # 슬라이드별로 청크 생성
                chunk_texts = self._split_text(full_text)
                for idx, chunk_text in enumerate(chunk_texts):
                    chunks.append(DocumentChunk(
                        text=chunk_text,
                        page=slide_num,
                        metadata={
                            "file_path": str(file_path),
                            "file_name": file_path.name,
                            "file_type": "pptx",
                            "slide": slide_num,
                            "total_slides": len(prs.slides),
                            "chunk_index": idx
                        }
                    ))
        
        logger.info(f"Parsed {len(chunks)} chunks from PPTX")
        return chunks
    
    def _parse_xlsx(self, file_path: Path) -> List[DocumentChunk]:
        """XLSX 파일 파싱 (openpyxl 사용)"""
        try:
            from openpyxl import load_workbook
        except ImportError:
            raise ImportError("openpyxl is required. Install: pip install openpyxl")
        
        chunks = []
        wb = load_workbook(file_path, data_only=True)
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            
            # 시트의 모든 셀 데이터 수집
            sheet_data = []
            for row in sheet.iter_rows(values_only=True):
                # 빈 행 제외
                row_text = [str(cell) for cell in row if cell is not None]
                if row_text:
                    sheet_data.append(" | ".join(row_text))
            
            # 시트 텍스트 결합
            full_text = "\n".join(sheet_data)
            
            if full_text.strip():
                # 시트별로 청크 생성
                chunk_texts = self._split_text(full_text)
                for idx, chunk_text in enumerate(chunk_texts):
                    chunks.append(DocumentChunk(
                        text=chunk_text,
                        section=sheet_name,
                        metadata={
                            "file_path": str(file_path),
                            "file_name": file_path.name,
                            "file_type": "xlsx",
                            "sheet_name": sheet_name,
                            "chunk_index": idx
                        }
                    ))
        
        logger.info(f"Parsed {len(chunks)} chunks from XLSX")
        return chunks
    
    def _parse_text(self, file_path: Path) -> List[DocumentChunk]:
        """TXT, MD 파일 파싱"""
        chunks = []
        
        # 인코딩 자동 감지
        encodings = ['utf-8', 'cp949', 'euc-kr']
        text = None
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    text = f.read()
                break
            except UnicodeDecodeError:
                continue
        
        if text is None:
            raise ValueError(f"Could not decode file: {file_path}")
        
        # 청크로 분할
        chunk_texts = self._split_text(text)
        for idx, chunk_text in enumerate(chunk_texts):
            chunks.append(DocumentChunk(
                text=chunk_text,
                metadata={
                    "file_path": str(file_path),
                    "file_name": file_path.name,
                    "file_type": file_path.suffix[1:],
                    "chunk_index": idx
                }
            ))
        
        logger.info(f"Parsed {len(chunks)} chunks from text file")
        return chunks
    
    def _split_text(self, text: str) -> List[str]:
        """
        텍스트를 청크로 분할
        
        Args:
            text: 분할할 텍스트
            
        Returns:
            청크 리스트
        """
        if not text.strip():
            return []
        
        chunks = []
        start = 0
        text_length = len(text)
        
        while start < text_length:
            end = start + self.chunk_size
            
            # 마지막 청크가 아니면 단어 경계에서 자르기
            if end < text_length:
                # 공백이나 줄바꿈 찾기
                while end > start and text[end] not in [' ', '\n', '\t', '.', ',', '!', '?']:
                    end -= 1
                
                # 적절한 분할점을 못 찾으면 그냥 자르기
                if end == start:
                    end = start + self.chunk_size
            
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            
            # 다음 청크 시작점 (오버랩 적용)
            start = end - self.chunk_overlap if end < text_length else text_length
        
        return chunks
    
    @classmethod
    def is_supported(cls, file_path: Path) -> bool:
        """파일이 지원되는 형식인지 확인"""
        return file_path.suffix.lower() in cls.SUPPORTED_EXTENSIONS

